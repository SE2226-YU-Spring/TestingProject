"""Build SE 2226 Test Project presentation following docs/presentation_rules.pptx structure."""
from copy import deepcopy
from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

SRC = "docs/Presentation.pptx"
OUT = "docs/Presentation.pptx"

prs = Presentation(SRC)

LAYOUT_TITLE = prs.slide_layouts[0]      # Başlık Slaydı
LAYOUT_CONTENT = prs.slide_layouts[1]    # Başlık ve İçerik
LAYOUT_SECTION = prs.slide_layouts[2]    # Bölüm Üstbilgisi
LAYOUT_TWO = prs.slide_layouts[3]        # İki İçerik
LAYOUT_TITLE_ONLY = prs.slide_layouts[5] # Yalnızca Başlık
LAYOUT_BLANK = prs.slide_layouts[6]      # Boş


def remove_slide(prs, idx):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    rId = slides[idx].get(qn("r:id"))
    prs.part.drop_rel(rId)
    xml_slides.remove(slides[idx])


def set_title(slide, title_text):
    if slide.shapes.title is not None:
        slide.shapes.title.text = title_text
        for p in slide.shapes.title.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True


def get_body_placeholder(slide):
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            return shape
    return None


def fill_bullets(slide, bullets, font_size=18):
    body = get_body_placeholder(slide)
    if body is None:
        return
    tf = body.text_frame
    tf.word_wrap = True
    tf.clear()
    first = True
    for item in bullets:
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = text
        p.level = level
        for run in p.runs:
            run.font.size = Pt(font_size)


def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = text
    for p in tf.paragraphs:
        p.alignment = align
        for r in p.runs:
            r.font.size = Pt(font_size)
            r.font.bold = bold
    return tb


def add_table(slide, left, top, width, height, data, header_color=RGBColor(0x2E, 0x5A, 0x88), header_font_size=12, body_font_size=11):
    rows = len(data)
    cols = len(data[0])
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = tbl_shape.table
    for r_idx, row in enumerate(data):
        for c_idx, cell_text in enumerate(row):
            cell = tbl.cell(r_idx, c_idx)
            cell.text = str(cell_text)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    if r_idx == 0:
                        run.font.bold = True
                        run.font.size = Pt(header_font_size)
                        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    else:
                        run.font.size = Pt(body_font_size)
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
    return tbl


# Remove the existing empty content slide (slide #2), keep cover (slide #1)
while len(prs.slides) > 1:
    remove_slide(prs, 1)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


# ------------------------------------------------------------------
# Slide 2 — CONTENT / OUTLINE
# ------------------------------------------------------------------
s = prs.slides.add_slide(LAYOUT_CONTENT)
set_title(s, "CONTENT / OUTLINE")
fill_bullets(s, [
    ("1. Introduction", 0),
    ("Project goal, scope, team", 1),
    ("2. Application Analysis", 0),
    ("Yemeksepeti (Black-Box) & Library-DB-Express (White-Box)", 1),
    ("3. Black-Box Testing", 0),
    ("Equivalence Partitioning, Boundary Value Analysis, Use Case, Decision Table", 1),
    ("Web (Playwright) and Mobile (Maestro) automation", 1),
    ("4. White-Box Testing", 0),
    ("Statement / Branch coverage with Jest", 1),
    ("Mutation Testing with Stryker", 1),
    ("5. Results & Lessons Learned", 0),
    ("6. Q&A / Contact", 0),
], font_size=18)


# ------------------------------------------------------------------
# Slide 3 — OVERVIEW (section header style)
# ------------------------------------------------------------------
s = prs.slides.add_slide(LAYOUT_SECTION)
set_title(s, "OVERVIEW")
body = get_body_placeholder(s)
if body:
    body.text_frame.text = (
        "Problem Definition  —  Application Analysis  —  "
        "Test Cases  —  Automated Tests"
    )
    for p in body.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(20)


# ------------------------------------------------------------------
# Slide 4 — Problem Definition
# ------------------------------------------------------------------
s = prs.slides.add_slide(LAYOUT_CONTENT)
set_title(s, "1. PROBLEM DEFINITION")
fill_bullets(s, [
    "Software Quality Assurance and Testing course project (SE 2226).",
    "Apply both Black-Box and White-Box testing techniques to real applications.",
    ("Black-Box target: Yemeksepeti — large food delivery platform (web + mobile).", 0),
    ("White-Box target: Library-DB-Express — Node.js / Express + Sequelize app.", 0),
    "Goals:",
    ("Cover Equivalence Partitioning, BVA, Use Case, Decision Table.", 1),
    ("Reach high statement / branch coverage and a strong mutation score.", 1),
    ("Produce automation that can be re-run in CI.", 1),
], font_size=18)


# ------------------------------------------------------------------
# Slide 5 — Application Analysis (two columns)
# ------------------------------------------------------------------
s = prs.slides.add_slide(LAYOUT_TWO)
set_title(s, "2. APPLICATION ANALYSIS")
# layout 3 has two body placeholders, idx 1 and 2
left_body = right_body = None
for ph in s.placeholders:
    if ph.placeholder_format.idx == 1:
        left_body = ph
    elif ph.placeholder_format.idx == 2:
        right_body = ph

def fill_two(ph, header, lines):
    tf = ph.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.text = header
    for r in p.runs:
        r.font.bold = True
        r.font.size = Pt(20)
    for line in lines:
        p = tf.add_paragraph()
        p.text = line
        p.level = 1
        for r in p.runs:
            r.font.size = Pt(15)

if left_body:
    fill_two(left_body, "Yemeksepeti (Black-Box)", [
        "Food delivery (web + mobile).",
        "Account, search, browse, cart, checkout.",
        "Web: Chrome via Playwright/Java.",
        "Mobile: Maestro on Android.",
        "Synthetic data only — no real payment.",
    ])
if right_body:
    fill_two(right_body, "Library-DB-Express (White-Box)", [
        "Node.js + Express + Sequelize + SQLite.",
        "Books CRUD, search, pagination.",
        "Routes, models, error handlers, server bootstrap.",
        "Tested with Jest 50 cases.",
        "Mutation tested with Stryker v7.3.0.",
    ])


# ------------------------------------------------------------------
# Slide 6 — Black-Box Testing section header
# ------------------------------------------------------------------
s = prs.slides.add_slide(LAYOUT_SECTION)
set_title(s, "3. BLACK-BOX TESTING")
body = get_body_placeholder(s)
if body:
    body.text_frame.text = "Yemeksepeti — Web (Playwright) and Mobile (Maestro)"
    for p in body.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(20)


# ------------------------------------------------------------------
# Slide 7 — Black-Box scope & techniques
# ------------------------------------------------------------------
s = prs.slides.add_slide(LAYOUT_CONTENT)
set_title(s, "BLACK-BOX SCOPE & TECHNIQUES")
fill_bullets(s, [
    "Test level: System Testing — Functional + UI + Regression (test_plan.pdf §1.1).",
    "Out of scope: real payment, performance / load / security (§1.3).",
    "Techniques applied (§3.5):",
    ("Equivalence Partitioning", 1),
    ("Boundary Value Analysis", 1),
    ("Use Case Testing", 1),
    ("Decision Table Testing", 1),
    "Tools:",
    ("Web: Java + Playwright 1.49 (Selenium-equivalent), JUnit 5.", 1),
    ("Mobile: Maestro Studio on a physical Android device.", 1),
], font_size=18)


# ------------------------------------------------------------------
# Slide 8 — Equivalence Partitioning example (Search)
# ------------------------------------------------------------------
s = prs.slides.add_slide(LAYOUT_TITLE_ONLY)
set_title(s, "EQUIVALENCE PARTITIONING — Restaurant Search")
add_text_box(
    s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.6),
    "Address / search query partitioned into equivalence classes.",
    font_size=16,
)
table_data = [
    ["Class", "Example input", "Expected", "Observed"],
    ["Valid full address", "\"Hatay\"", "Suggestions list", "PASS"],
    ["Valid partial token", "\"Üniversite 2\"", "Filtered suggestions", "PASS"],
    ["Non-existent string", "\"RandomText\"", "\"No results\" empty state", "PASS"],
    ["Empty query", "\"\"", "Submit disabled / no call", "PASS"],
    ["Special characters", "\"!@#$%\"", "No crash, empty state", "PASS"],
]
add_table(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(4.3), table_data)


# ------------------------------------------------------------------
# Slide 9 — Boundary Value Analysis (Cart quantity)
# ------------------------------------------------------------------
s = prs.slides.add_slide(LAYOUT_TITLE_ONLY)
set_title(s, "BOUNDARY VALUE ANALYSIS — Cart Quantity")
add_text_box(
    s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.6),
    "Sequence 0 → 1 → 2 → 1 → 0 verified on web (cart) and mobile.",
    font_size=16,
)
table_data = [
    ["Equivalence class", "Input", "Expected result", "Actual result"],
    ["U1 — Below lower bound", "0 (remove)", "Item removed, cart empty", "Cart empty"],
    ["E1 — Lower bound", "1", "Item added, total = price × 1", "Match"],
    ["E2 — Increment", "2", "Total = price × 2", "Match"],
    ["E3 — Decrement back to 1", "1", "Total restored to price × 1", "Match"],
    ["U2 — Remove from 1", "0", "Cart returns to empty", "Match"],
]
add_table(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(4.3), table_data)


# ------------------------------------------------------------------
# Slide 10 — Decision Table (Checkout payment)
# ------------------------------------------------------------------
s = prs.slides.add_slide(LAYOUT_TITLE_ONLY)
set_title(s, "DECISION TABLE — Checkout Payment Method")
add_text_box(
    s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.6),
    "Each rule is an input combination; expected outcome is the order placement state.",
    font_size=16,
)
table_data = [
    ["Rule", "Cart filled", "Address selected", "Payment", "Expected"],
    ["R1", "Yes", "Yes", "Cash (Nakit)", "Place-order enabled"],
    ["R2", "Yes", "Yes", "Online card", "Place-order enabled"],
    ["R3", "Yes", "Yes", "Card on delivery", "Place-order enabled"],
    ["R4", "Yes", "No", "any", "Address step required"],
    ["R5", "No", "any", "any", "Checkout disabled"],
]
add_table(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(4.3), table_data)


# ------------------------------------------------------------------
# Slide 11 — Use Case Testing (Add to cart use case)
# ------------------------------------------------------------------
s = prs.slides.add_slide(LAYOUT_CONTENT)
set_title(s, "USE CASE TESTING — Add Product to Cart")
fill_bullets(s, [
    "Actors: Guest user (web), authenticated user (mobile).",
    "Pre-conditions: app reachable, address picked, restaurant open.",
    "Main flow:",
    ("1. Open menu → 2. Pick item → 3. Add to cart", 1),
    ("4. Increment quantity → 5. Verify total math", 1),
    ("6. Decrement → 7. Remove → 8. Cart empty", 1),
    "Post-conditions: cart state matches user actions; price math correct.",
    "Result: ✅ All steps pass for web Playwright run on 2026-05-07.",
], font_size=18)


# ------------------------------------------------------------------
# Slide 12 — Mobile (Maestro) tests
# ------------------------------------------------------------------
s = prs.slides.add_slide(LAYOUT_CONTENT)
set_title(s, "MOBILE TESTS — Maestro Studio")
fill_bullets(s, [
    "Three flows authored as Maestro YAML scripts:",
    ("YemekSepetiSearch.yaml — address & restaurant search", 1),
    ("YemekSepetiCartTest.yaml — add / increment / decrement / remove", 1),
    ("YemekSepetiPayTest.yaml — address book + payment selectors", 1),
    "Run on a physical Android device via ADB.",
    "Selectors anchored on Maestro element IDs (HomeSearchBar, AUTOCOMPLETE_SUGGESTION_ENTRY).",
    "Assertions cover navigation, visible text, and cart totals.",
], font_size=18)


# ------------------------------------------------------------------
# Slide 13 — Web (Playwright) tests + results
# ------------------------------------------------------------------
s = prs.slides.add_slide(LAYOUT_TITLE_ONLY)
set_title(s, "WEB TESTS — Playwright/JUnit Results")
add_text_box(
    s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.6),
    "Run on 2026-05-07, Chrome 148, persistent profile, locale tr-TR.",
    font_size=16,
)
table_data = [
    ["#", "Test class", "Method", "Result", "Time"],
    ["1", "YemekSepetiSearchTest", "scenario1_addressSuggestionAndRestaurantDetail", "PASS", "~26 s"],
    ["2", "YemekSepetiSearchTest", "scenario2_topbarSearchAndDetail", "PASS", "~28 s"],
    ["3", "YemekSepetiPayTest", "scenario1_addressModalNewAddress", "PASS", "~22 s"],
    ["4", "YemekSepetiPayTest", "scenario2_paymentMethodFiltersAreSelectable", "PASS", "~23 s"],
    ["5", "YemekSepetiCartTest", "cartAddIncrementDecrementRemove", "PASS", "~59 s"],
]
add_table(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(3.5), table_data)
add_text_box(
    s, Inches(0.5), Inches(5.7), Inches(12.3), Inches(0.6),
    "Aggregate: 5 / 5 passing — 0 logic failures.",
    font_size=16, bold=True,
)


# ------------------------------------------------------------------
# Slide 14 — Black-Box findings & gaps
# ------------------------------------------------------------------
s = prs.slides.add_slide(LAYOUT_CONTENT)
set_title(s, "BLACK-BOX FINDINGS & GAPS")
fill_bullets(s, [
    "Defects observed: none on the happy path during the recorded run.",
    "Gaps:",
    ("Mobile address-book editing partially behind login on web — verified at API-id level only.", 1),
    ("Anti-bot (PerimeterX) interferes with unattended runs; mitigated via persistent profile pre-warm.", 1),
    "Test data is synthetic only (test_plan.pdf §3.9).",
    "Cart assertions are price-agnostic — verify the math, not literal TL strings.",
], font_size=18)


# ------------------------------------------------------------------
# Slide 15 — White-Box section header
# ------------------------------------------------------------------
s = prs.slides.add_slide(LAYOUT_SECTION)
set_title(s, "4. WHITE-BOX TESTING")
body = get_body_placeholder(s)
if body:
    body.text_frame.text = "Library-DB-Express — Jest unit + integration, Stryker mutation"
    for p in body.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(20)


# ------------------------------------------------------------------
# Slide 16 — White-Box scope
# ------------------------------------------------------------------
s = prs.slides.add_slide(LAYOUT_CONTENT)
set_title(s, "WHITE-BOX SCOPE")
fill_bullets(s, [
    "Test levels: Unit + Integration.",
    "Source under test:",
    ("routes/index.js — CRUD + pagination + search", 1),
    ("models/book.js — Sequelize validators", 1),
    ("errorHandlers.js — 404 and global error", 1),
    ("app.js + bin/www — middleware order, server bootstrap", 1),
    "Test framework: Jest with in-memory SQLite for isolation.",
    "Mutation framework: Stryker v7.3.0 (Node.js v18 LTS).",
], font_size=18)


# ------------------------------------------------------------------
# Slide 17 — Coverage table
# ------------------------------------------------------------------
s = prs.slides.add_slide(LAYOUT_TITLE_ONLY)
set_title(s, "WHITE-BOX COVERAGE")
table_data = [
    ["Component", "Line coverage", "Branch coverage", "Tests"],
    ["routes/index.js", "92%", "85%", "18"],
    ["models/book.js", "95%", "90%", "8"],
    ["errorHandlers.js", "100%", "95%", "6"],
    ["app.js", "88%", "75%", "5"],
    ["mutation.test.js (extra robustness)", "—", "—", "13"],
    ["Overall", "~92%", "~85%", "50"],
]
add_table(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(4.5), table_data)
add_text_box(
    s, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.6),
    "All 50 Jest tests pass — 0 failures, 0 skipped.",
    font_size=16, bold=True,
)


# ------------------------------------------------------------------
# Slide 18 — Mutation Testing
# ------------------------------------------------------------------
s = prs.slides.add_slide(LAYOUT_TITLE_ONLY)
set_title(s, "MUTATION TESTING — Stryker v7.3.0")
add_text_box(
    s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.6),
    "Mutation Score = killed / (killed + survived) = 145 / 160 = 90.67%",
    font_size=18, bold=True,
)
table_data = [
    ["Metric", "Count", "%"],
    ["Mutants generated", "164", "100%"],
    ["Killed (caught by tests)", "145", "88.4%"],
    ["Survived", "14", "8.5%"],
    ["Timed out", "5", "3.0%"],
    ["Mutation score (valid mutants)", "145 / 160", "90.67%"],
]
add_table(s, Inches(0.5), Inches(2.1), Inches(12.3), Inches(3.6), table_data)
add_text_box(
    s, Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.6),
    "Threshold target ≥ 80% — exceeded by 10.67 percentage points.",
    font_size=16,
)


# ------------------------------------------------------------------
# Slide 19 — Survived mutants categories
# ------------------------------------------------------------------
s = prs.slides.add_slide(LAYOUT_CONTENT)
set_title(s, "SURVIVED MUTANTS — Categories")
fill_bullets(s, [
    "Boolean operator mutations (5–6) — e.g. && → || , > → >=",
    "Optional / edge cases (3–4) — null vs undefined, whitespace vs empty",
    "Default-value mutations (2–3) — default parameter / constant changes",
    "Timeout cases (5) — 60 s budget exceeded; possible long async loops",
    "Action plan:",
    ("Add boundary value tests (0, -1, empty, whitespace-only)", 1),
    ("Add boolean-combination tests around pagination math", 1),
    ("Raise timeout to 120 s for long async tests", 1),
], font_size=18)


# ------------------------------------------------------------------
# Slide 20 — Conclusions / Future Work
# ------------------------------------------------------------------
s = prs.slides.add_slide(LAYOUT_CONTENT)
set_title(s, "5. CONCLUSIONS & FUTURE WORK")
fill_bullets(s, [
    "Black-Box (Yemeksepeti):",
    ("All five web scenarios pass; mobile flows automated in Maestro.", 1),
    ("Authenticated checkout coverage gated by login — flagged as a gap.", 1),
    "White-Box (Library-DB-Express):",
    ("50 / 50 Jest tests passing; ~92% line / ~85% branch coverage.", 1),
    ("Mutation score 90.67% — well above the 80% threshold.", 1),
    "Future work:",
    ("Add a synthetic test account for full checkout-side assertions.", 1),
    ("Push mutation score to ≥ 95% with targeted boundary tests.", 1),
    ("Wire Jest + Stryker into CI with score thresholds.", 1),
], font_size=17)


# ------------------------------------------------------------------
# Slide 21 — Lessons Learned
# ------------------------------------------------------------------
s = prs.slides.add_slide(LAYOUT_CONTENT)
set_title(s, "LESSONS LEARNED")
fill_bullets(s, [
    "Whitebox tests must mirror the actual source — assumption-driven tests caused rework.",
    "High line coverage ≠ high mutation detection — boundary tests are essential.",
    "Integration tests with a real DB caught issues that mocks would have hidden.",
    "Tool-version compatibility matters — Stryker v8 needed Node ≥ 20; we pinned v7.3.0 for Node 18.",
    "For black-box web automation, data-testid selectors are far more durable than CSS classes.",
], font_size=18)


# ------------------------------------------------------------------
# Slide 22 — Thank You
# ------------------------------------------------------------------
s = prs.slides.add_slide(LAYOUT_TITLE)
# title slide layout: idx 0 center title, idx 1 subtitle
title = s.shapes.title
title.text = "THANK YOU FOR LISTENING"
for p in title.text_frame.paragraphs:
    for r in p.runs:
        r.font.size = Pt(54)
        r.font.bold = True
sub = None
for ph in s.placeholders:
    if ph.placeholder_format.idx == 1:
        sub = ph
if sub:
    sub.text_frame.text = "Questions & Discussion"
    for p in sub.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(24)


# ------------------------------------------------------------------
# Slide 23 — Contact info (additional)
# ------------------------------------------------------------------
s = prs.slides.add_slide(LAYOUT_CONTENT)
set_title(s, "CONTACT INFO")
fill_bullets(s, [
    "Batuhan Yerebasmaz — 23070006029",
    "Osman Şahin Güler — 24070006009 — osmansahinguler@gmail.com",
    "İsmet Saygın Koç — 23070006038",
    "Ege Çınar — 24070006020",
    "Atakan Sezginer — 24070006126",
    "",
    "Course: SE 2226 — Software Quality Assurance and Testing",
    "Department of Software Engineering — Doç. Dr. Korhan Karabulut",
], font_size=18)


prs.save(OUT)
print(f"Wrote {OUT} with {len(prs.slides)} slides.")
